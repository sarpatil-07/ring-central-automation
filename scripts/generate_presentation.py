#!/usr/bin/env python3
"""Generate RingCentralAutoSet presentation (.pptx).

Run from project root:
  .venv/bin/python3.14 scripts/generate_presentation.py
  # or: .venv/bin/python scripts/generate_presentation.py  (if venv python matches pip)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "RingCentralAutoSet_MAX_Automation.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    box2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(8.8), Inches(1.5))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(5.0), Inches(1.2), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def _add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, ORANGE)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)

    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.7))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    top = 1.25 if not subtitle else 1.45
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(8.9), Inches(0.4))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = GRAY
        sp.font.italic = True

    body = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("  "):
            p.text = line.strip()
            p.level = 1
            p.font.size = Pt(15)
        elif line.startswith("|") or line.startswith("```"):
            p.text = line
            p.font.size = Pt(11)
            p.font.name = "Courier New"
        elif line == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            p.text = line
            p.level = 0
            p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.7))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = WHITE

    cols = len(headers)
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.45 * (len(rows) + 2)))
    table = table_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF1)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "RingCentralAutoSet",
        "Automated MAX Agent Login, Status & Logout on Schedule\n"
        "Browser-based Ring Central / NICE MAX automation for macOS",
    )

    _add_content_slide(
        prs,
        "The Problem",
        [
            "Agents repeat the same steps every shift:",
            "  • Shift start — Open MAX, login/OTP, connect Station ID, set Available",
            "  • Lunch — Change status to Unavailable (Lunch)",
            "  • After lunch — Set Available again",
            "  • Shift end — Logout / disconnect",
            "",
            "Pain points:",
            "  • Easy to forget login or logout at shift boundaries",
            "  • Repetitive manual clicks every day",
            "  • Work Chrome and MAX Chrome get mixed together",
        ],
    )

    _add_content_slide(
        prs,
        "Our Solution — RingCentralAutoSet",
        [
            "Configure once → runs automatically on schedule",
            "",
            "  • Opens a dedicated MAX Chrome window (not your work browser)",
            "  • Connects with your Station ID",
            "  • Sets Available / Lunch / Logout on schedule",
            "  • Runs in background after Mac login (optional)",
            "  • Supports IST, APAC, rotating shifts, and leave days",
            "",
            "Project: Ring_Central_Automation  |  App name: RingCentralAutoSet",
        ],
    )

    _add_content_slide(
        prs,
        "High-Level Architecture",
        [
            "macOS LaunchAgent (optional)",
            "  → APScheduler (cron by timezone & work days)",
            "  → run.py (CLI + scheduler)",
            "  → Playwright",
            "  → Dedicated Chrome (chrome-max-profile)",
            "  → NICE MAX web UI automation",
            "",
            "Everything runs locally on the agent's Mac.",
            "Schedule and Station ID stored in .env (not committed to git).",
        ],
    )

    _add_table_slide(
        prs,
        "Daily Automation Flow",
        ["Time", "Action", "What happens"],
        [
            ["Work start", "morning", "Open MAX Chrome → connect Station ID → Available"],
            ["Lunch start", "lunch", "Set Unavailable (Lunch) — quiet mode"],
            ["Lunch end", "lunch-end", "Set Available"],
            ["Work end", "logout", "Logout + close MAX browser"],
        ],
    )

    _add_table_slide(
        prs,
        "Why a Separate Chrome Profile?",
        ["", "Work Chrome", "MAX Chrome"],
        [
            ["Profile", "Your daily profile", "chrome-max-profile"],
            ["Purpose", "Email, tools, etc.", "MAX agent only"],
            ["Touched by script?", "Never", "Yes"],
            ["SSO session", "Your work login", "Saved MAX login"],
        ],
    )

    _add_table_slide(
        prs,
        "Technology Stack",
        ["Layer", "Technology", "Role"],
        [
            ["Language", "Python 3", "Main automation logic"],
            ["Browser", "Playwright", "Control MAX web UI"],
            ["Scheduler", "APScheduler", "Cron jobs by timezone"],
            ["Config", ".env + YAML", "Times, Station ID, selectors"],
            ["Background", "macOS LaunchAgent", "Auto-start on login"],
            ["Timezone", "pytz / IANA", "IST, APAC, EMEA, etc."],
        ],
    )

    _add_content_slide(
        prs,
        "Key Project Files",
        [
            "run.py              — Main entry point & scheduler",
            "max_flow.py         — Morning / lunch / logout workflows",
            "browser_session.py  — MAX Chrome launch & reuse",
            "ui_actions.py       — UI clicks & smart status checks",
            "schedule_menu.py    — Interactive configuration menu",
            "config.py           — Settings, timezone, schedule",
            "service.py          — macOS LaunchAgent installer",
            "selectors.yaml      — MAX UI element selectors",
            ".env                — Station ID, work times, timezone",
        ],
    )

    _add_table_slide(
        prs,
        "Default Schedule (Configurable)",
        ["Job", "Default", "Action"],
        [
            ["Work start", "09:00", "Connect + Available"],
            ["Lunch start", "13:00", "Unavailable (Lunch)"],
            ["Lunch end", "14:00", "Available"],
            ["Work end", "18:00", "Logout + close browser"],
        ],
    )

    _add_content_slide(
        prs,
        "Smart Features",
        [
            "Quiet lunch/logout — status jobs don't steal focus from work Chrome",
            "Leave / pause — skip jobs on vacation or mark a leave day",
            "Rotating shifts — change times in menu; scheduler picks up changes",
            "Smart status — skip if you already changed status manually",
            "Browser recovery — reset stuck MAX Chrome without losing login",
            "Session reuse — one MAX tab for the whole workday",
            "Misfire grace — jobs still run if Mac was asleep (up to 1 hour late)",
        ],
    )

    _add_table_slide(
        prs,
        "Interactive Menu (No Coding Required)",
        ["Key", "Function"],
        [
            ["1", "Show schedule"],
            ["2", "Set all times (work + lunch)"],
            ["8", "Change Station ID"],
            ["9", "Change timezone"],
            ["p / u", "Pause / resume automation"],
            ["x", "Mark today as leave + logout"],
            ["i", "Install background service"],
            ["r", "Reset stuck MAX browser"],
        ],
    )

    _add_content_slide(
        prs,
        "One-Time Setup",
        [
            "1. Copy project to ~/Projects/Ring_Central_Automation",
            "2. Run: bash setup.sh  (creates venv, installs Playwright)",
            "3. Run: .venv/bin/python run.py menu",
            "4. Set Station ID and shift times (menu options 8 and 2)",
            "5. Install background service: menu i  or  install-service",
            "6. First morning: complete OTP in MAX Chrome when prompted",
            "",
            "Requirements: macOS, Python 3, Google Chrome, network to MAX/CXone",
        ],
    )

    _add_section_slide(prs, "Three Automation Paths")

    _add_table_slide(
        prs,
        "Automation Options Built",
        ["Path", "Script", "Best for", "IT API?"],
        [
            ["Browser MAX", "run.py", "MAX in Chrome — recommended", "No"],
            ["Desktop app", "desktop_run.py", "RingCentral.app (Accessibility)", "No"],
            ["Official API", "rc_run.py", "Production desktop MAX", "Yes"],
        ],
    )

    _add_content_slide(
        prs,
        "Browser Path — How It Works (Detail)",
        [
            "Morning:",
            "  1. Launch or reuse MAX Chrome window",
            "  2. Wait for manual SSO/OTP if needed (script does not click during login)",
            "  3. Open MAX agent application from myprofile",
            "  4. Connect with Station ID",
            "  5. Set status → Available",
            "",
            "Lunch / logout reuse the same MAX tab for the whole workday.",
        ],
    )

    _add_content_slide(
        prs,
        "Security & Safety",
        [
            "Station ID stored in local .env file (gitignored)",
            "SSO / OTP handled only in dedicated Chrome profile",
            "Work browser is never automated",
            "No data sent to external servers — local Mac automation only",
            "Logs stored locally under logs/",
            "",
            "Recommendation: Confirm with team policy before wide rollout.",
        ],
    )

    _add_table_slide(
        prs,
        "Troubleshooting",
        ["Issue", "Fix"],
        [
            ["MAX Chrome stuck", "run.py reset-browser  or  menu r"],
            ["OTP required again", "run.py clear-session  or  menu k"],
            ["Jobs not running", "run.py service-status → reinstall service"],
            ["Wrong timezone", "Set TIMEZONE=Asia/Kolkata in .env"],
            ["Check logs", "logs/scheduler.err.log"],
        ],
    )

    _add_content_slide(
        prs,
        "Demo Commands",
        [
            "cd ~/Projects/Ring_Central_Automation",
            "",
            ".venv/bin/python run.py show        # Show schedule",
            ".venv/bin/python run.py morning     # Run work start now",
            ".venv/bin/python run.py menu          # Interactive menu",
            ".venv/bin/python run.py install-service  # Auto on Mac login",
        ],
    )

    _add_content_slide(
        prs,
        "Summary",
        [
            "RingCentralAutoSet automates the MAX agent daily routine on Mac",
            "",
            "  • Separate MAX Chrome — safe from work browser",
            "  • Scheduled Available / Lunch / Logout",
            "  • Menu-driven config — no code for shift changes",
            "  • Background service for hands-free daily use",
            "",
            "Project folder: ~/Projects/Ring_Central_Automation",
            "",
            "Questions?",
        ],
    )

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
