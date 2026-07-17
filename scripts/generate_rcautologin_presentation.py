#!/usr/bin/env python3
"""Generate RCAutoLogin team demo presentation (.pptx).

Run from project root:
  .venv/bin/python -m pip install python-pptx   # once
  .venv/bin/python scripts/generate_rcautologin_presentation.py
"""

from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError:
    print("Missing dependency: python-pptx")
    print("Install once, then re-run this script:")
    print("  .venv/bin/python -m pip install -r requirements-dev.txt")
    print("  .venv/bin/python scripts/generate_rcautologin_presentation.py")
    raise SystemExit(1)

OUT = Path(__file__).resolve().parent.parent / "RCAutoLogin_Team_Demo.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x06, 0x67, 0xC2)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x23, 0x32)
GRAY = RGBColor(0x5C, 0x6B, 0x7A)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
SOFT_BLUE = RGBColor(0xE8, 0xF2, 0xFC)
GREEN = RGBColor(0x0D, 0x7A, 0x3F)
PLACEHOLDER = RGBColor(0xD0, 0xD7, 0xDE)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.7))
    p = tbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    box2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.15), Inches(8.8), Inches(1.6))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.0), Inches(1.2), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def _add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, ORANGE)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def _add_bullets(slide, top: float, bullets: list[str], *, width: float = 8.9, size: int = 18) -> None:
    body = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("  "):
            p.text = line.strip()
            p.level = 1
            p.font.size = Pt(size - 2)
        elif line == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            p.text = line
            p.level = 0
            p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, title)
    _add_bullets(slide, 1.25, bullets)


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, title)
    cols = len(headers)
    height = min(0.45 * (len(rows) + 2), 4.8)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.4), Inches(1.35), Inches(9.2), Inches(height)
    )
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF1)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK


def _box(slide, left, top, width, height, text: str, *, fill: RGBColor = SOFT_BLUE, font_size: int = 13) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER


def _arrow_down(slide, x, y1, y2) -> None:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    conn.line.color.rgb = NAVY
    conn.line.width = Pt(2)


def _add_architecture_diagram(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, "Architecture — How It Works")

    # Top row: User controls
    _box(slide, Inches(0.5), Inches(1.4), Inches(2.8), Inches(0.85), "Web GUI\n(browser tab)", fill=SOFT_BLUE)
    _box(slide, Inches(3.6), Inches(1.4), Inches(2.8), Inches(0.85), "Background\nscheduler", fill=RGBColor(0xFF, 0xF4, 0xE5))

    _arrow_down(slide, Inches(1.9), Inches(2.25), Inches(2.65))
    _arrow_down(slide, Inches(5.0), Inches(2.25), Inches(2.65))

    # Center: RCAutoLogin
    _box(
        slide,
        Inches(2.5),
        Inches(2.65),
        Inches(5.0),
        Inches(0.9),
        "RCAutoLogin  (Python + Playwright)",
        fill=RGBColor(0xDC, 0xF5, 0xE7),
        font_size=15,
    )

    _arrow_down(slide, Inches(5.0), Inches(3.55), Inches(3.95))

    # Chrome
    _box(
        slide,
        Inches(1.8),
        Inches(3.95),
        Inches(6.4),
        Inches(0.85),
        "RingCX Chrome  (separate profile — not work Chrome)",
        fill=WHITE,
    )

    _arrow_down(slide, Inches(5.0), Inches(4.8), Inches(5.2))

    # RingCX cloud
    _box(
        slide,
        Inches(1.8),
        Inches(5.2),
        Inches(6.4),
        Inches(0.85),
        "app.ringcentral.com  —  RingCX Agent",
        fill=NAVY,
        font_size=14,
    )
    tf = slide.shapes[-1].text_frame
    tf.paragraphs[0].font.color.rgb = WHITE

    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    p = note.text_frame.paragraphs[0]
    p.text = "Everything runs on the agent's laptop. No cloud server. GUI is optional after setup."
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.font.italic = True


def _add_setup_flow_diagram(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, "One-Time Setup (Each User)")

    steps = [
        ("1", "SETUP tab\nSave login\n(email + password)"),
        ("2", "SCHEDULE tab\nWork & lunch times"),
        ("3", "TODAY tab\nStart auto job"),
    ]
    x = 0.55
    for i, (num, label) in enumerate(steps):
        left = Inches(x + i * 3.1)
        _box(slide, left, Inches(1.6), Inches(2.7), Inches(1.3), f"Step {num}\n{label}", font_size=12)
        if i < 2:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.75), Inches(2.05), Inches(0.35), Inches(0.35))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    _box(slide, Inches(1.0), Inches(3.3), Inches(8.0), Inches(0.75), "Done — daily jobs run in background (GUI can be closed)", fill=RGBColor(0xDC, 0xF5, 0xE7))

    _add_bullets(
        slide,
        4.3,
        [
            "Login saved once in local .env — not shared in zip file",
            "Chrome stays open after manual Login; only Logout closes browser",
            "Mac: LaunchAgent  |  Linux: systemd user service",
        ],
        width=8.9,
        size=16,
    )


def _add_daily_flow_diagram(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, "Daily Automation Flow")

    flows = [
        ("Work start", "Login → Agent → Start session → AVAILABLE", GREEN),
        ("Lunch", "Set LUNCH status", ORANGE),
        ("After lunch", "Set AVAILABLE", GREEN),
        ("Work end", "Stop session → close browser", RGBColor(0xB4, 0x23, 0x18)),
    ]
    y = 1.45
    for title, desc, color in flows:
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(1.6), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(y), Inches(7.3), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.color.rgb = NAVY
        tf2 = bar.text_frame
        tf2.paragraphs[0].text = desc
        tf2.paragraphs[0].font.size = Pt(14)
        tf2.paragraphs[0].font.color.rgb = DARK
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        y += 0.75

    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.2))
    tf = note.text_frame
    tf.paragraphs[0].text = "Optional auto sign-in: saved credentials used when not already logged in."
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = GRAY
    p2 = tf.add_paragraph()
    p2.text = "MFA/OTP: complete once in RingCX Chrome — tool waits and continues."
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY


def _add_gui_tabs_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, "Web GUI — Four Tabs")

    tabs = [
        ("TODAY", "Login · Lunch · Back · Logout\nSchedule summary · Start auto job"),
        ("SETUP", "Save RingCentral login once\nStep 1-2-3 progress"),
        ("SCHEDULE", "Work start/end · Lunch · Timezone"),
        ("LOG", "Activity log · errors"),
    ]
    for i, (name, desc) in enumerate(tabs):
        row, col = divmod(i, 2)
        left = Inches(0.5 + col * 4.7)
        top = Inches(1.45 + row * 2.0)
        _box(slide, left, top, Inches(4.3), Inches(1.6), f"{name}\n\n{desc}", font_size=12)

    _add_bullets(slide, 5.5, ["Launch: .venv/bin/python rc_autologin_run.py  →  http://127.0.0.1:8765/"], size=14)


def _add_screenshot_slide(prs: Presentation, title: str, hint: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _title_bar(slide, title)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.3), Inches(8.9), Inches(5.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = PLACEHOLDER
    rect.line.color.rgb = GRAY
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Paste screenshot here"
    p.font.size = Pt(26)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = hint
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK
    p2.alignment = PP_ALIGN.CENTER


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "RCAutoLogin",
        "Automate RingCX Agent — Login, Lunch & Logout\n"
        "Local tool · Mac & Linux · Web GUI control panel",
    )

    _add_content_slide(
        prs,
        "The Problem",
        [
            "Every shift, RingCX web agents repeat the same steps:",
            "  • Sign in → Agent tab → Start session → AVAILABLE",
            "  • Lunch → LUNCH status → back to AVAILABLE",
            "  • End of day → Stop session",
            "",
            "Pain points:",
            "  • Easy to forget login or logout",
            "  • Same clicks every day",
            "  • RingCX mixed with normal work Chrome",
        ],
    )

    _add_content_slide(
        prs,
        "The Solution — RCAutoLogin",
        [
            "Runs on a schedule you set (work days + timezone)",
            "",
            "  • Opens RingCX in a dedicated Chrome window",
            "  • Work start → Login → AVAILABLE",
            "  • Lunch → LUNCH → AVAILABLE",
            "  • Work end → Stop session",
            "",
            "Simple web GUI — four tabs: Today, Setup, Schedule, Log",
            "Background job runs daily — even if GUI is closed",
        ],
    )

    _add_content_slide(
        prs,
        "What It Does NOT Do (Be Honest)",
        [
            "Not an IT API — it clicks the RingCX web UI like you do manually",
            "Needs a visible desktop (Mac/Linux + Google Chrome)",
            "MFA/OTP: you may approve once; tool waits and continues",
            "",
            "Optional auto sign-in:",
            "  • Save email + password once in SETUP tab",
            "  • Stored only in local .env on your machine",
            "  • Never included in the shared zip file",
        ],
    )

    _add_table_slide(
        prs,
        "Example Daily Schedule",
        ["Time", "What happens"],
        [
            ["09:00", "Login + Start session + AVAILABLE"],
            ["13:00", "LUNCH"],
            ["14:00", "AVAILABLE"],
            ["18:00", "Stop session + close browser"],
        ],
    )

    _add_architecture_diagram(prs)

    _add_setup_flow_diagram(prs)

    _add_daily_flow_diagram(prs)

    _add_table_slide(
        prs,
        "Work Chrome vs RingCX Chrome",
        ["", "Your work Chrome", "RingCX Chrome (RCAutoLogin)"],
        [
            ["Profile folder", "Default Chrome", "chrome-rcx-profile"],
            ["Touched by tool?", "Never", "Yes"],
            ["Login session", "Your work SSO", "RingCX SSO only"],
        ],
    )

    _add_gui_tabs_slide(prs)

    _add_section_slide(prs, "Live Demo & Screenshots")

    _add_screenshot_slide(prs, "Screenshot — TODAY tab", "Quick actions + Start auto job + schedule summary")
    _add_screenshot_slide(prs, "Screenshot — SETUP tab", "Save login once (email + hidden password)")
    _add_screenshot_slide(prs, "Screenshot — RingCX Chrome", "Agent tab · AVAILABLE status")

    _add_table_slide(
        prs,
        "Share Zip With Team",
        ["Who", "Steps"],
        [
            ["Builder", "build-release → dist/RCAutoLogin-1.0.0-portable.zip"],
            ["All users", "Unzip → cd RCAutoLogin-1.0.0-portable → ./install.sh"],
            ["All users", "GUI: Setup → Schedule → Start auto job"],
            ["Never share", ".env, chrome-rcx-profile, logs"],
        ],
    )

    _add_table_slide(
        prs,
        "Launch GUI — Mac vs Linux (After install.sh)",
        ["Step", "Mac", "Linux"],
        [
            ["1 — Install (once)", "./install.sh", "./install.sh  (same)"],
            ["2 — Open GUI", "Double-click Launch RCAutoLogin.command", "./launch-gui.sh"],
            ["Optional", "Or open RCAutoLogin.app from Dock", "Menu: rcautologin.desktop"],
            ["Browser opens", "http://127.0.0.1:8765/", "http://127.0.0.1:8765/  (same)"],
            ["Background job", "LaunchAgent (auto on login)", "systemd user service"],
            ["Requires", "Python 3.11–3.13, Google Chrome", "Python 3.11–3.13, Chrome/Chromium"],
        ],
    )

    _add_content_slide(
        prs,
        "Live Demo Script (15–20 min)",
        [
            "1. Show GUI TODAY tab — status pills",
            "2. SETUP tab — save login (one-time)",
            "3. SCHEDULE tab — work/lunch times",
            "4. TODAY → Login — Chrome stays open",
            "5. Start auto job — runs without GUI",
            "6. LOG tab — activity",
            "",
            "Full script: RCAutoLogin_DEMO_SCRIPT.txt",
            "Full guide: RCAutoLogin_COMPLETE_GUIDE.txt",
        ],
    )

    _add_title_slide(prs, "Questions?", "RCAutoLogin — RingCX web agent automation")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
